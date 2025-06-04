"""
Document Processor Module

This module is responsible for processing various document formats including
PDF, DOCX, CSV, PPTX, and Excel files with complete functionality.

Technologies: PyMuPDF, python-docx, pandas, python-pptx, pdfplumber
"""

import os
import time
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Optional, Union
import logging

# Import document processing libraries
try:
    import fitz  # PyMuPDF
    import docx
    import pandas as pd
    import pptx
    import pdfplumber
    from openpyxl import load_workbook
except ImportError as e:
    logging.warning(f"Some document processing libraries are not installed: {e}")

from utils.error_handler import DocumentProcessingError, error_handler, ErrorType


class DocumentProcessor:
    """
    Processes various document formats and extracts text content with full functionality.

    Supported formats:
    - PDF (using PyMuPDF and pdfplumber)
    - DOCX (using python-docx)
    - CSV/Excel (using pandas)
    - PPTX (using python-pptx)
    """

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """
        Initialize the DocumentProcessor with configuration.

        Args:
            config: Configuration dictionary with processing parameters
        """
        self.config = config or {}
        self.logger = logging.getLogger(__name__)

        # Configuration settings
        self.max_file_size_mb = self.config.get("max_file_size_mb", 50)
        self.supported_formats = self.config.get(
            "supported_formats",
            [".pdf", ".docx", ".csv", ".xlsx", ".xls", ".pptx", ".txt", ".md"],
        )

    @error_handler(ErrorType.DOCUMENT_PROCESSING)
    def process_document(self, file_path: str) -> Dict[str, Any]:
        """
        Process a document and extract its text content with metadata.

        Args:
            file_path: Path to the document file

        Returns:
            Dictionary containing extracted text and metadata
        """
        if not os.path.exists(file_path):
            raise DocumentProcessingError(f"Document not found: {file_path}", file_path)

        # Validate file size
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        if file_size_mb > self.max_file_size_mb:
            raise DocumentProcessingError(
                f"File too large: {file_size_mb:.1f}MB (max: {self.max_file_size_mb}MB)",
                file_path,
            )

        file_extension = os.path.splitext(file_path)[1].lower()

        # Validate file format
        if file_extension not in self.supported_formats:
            raise DocumentProcessingError(
                f"Unsupported file format: {file_extension}", file_path
            )

        self.logger.info(f"Processing document: {file_path} ({file_size_mb:.1f}MB)")

        try:
            if file_extension == ".pdf":
                return self._process_pdf(file_path)
            elif file_extension == ".docx":
                return self._process_docx(file_path)
            elif file_extension in [".csv", ".xlsx", ".xls"]:
                return self._process_spreadsheet(file_path)
            elif file_extension == ".pptx":
                return self._process_pptx(file_path)
            elif file_extension in [".txt", ".md"]:
                return self._process_text_file(file_path)
        except Exception as e:
            raise DocumentProcessingError(
                f"Error processing document: {str(e)}", file_path
            )

    def process_batch(self, file_paths: List[str]) -> List[Dict[str, Any]]:
        """
        Process multiple documents in batch.

        Args:
            file_paths: List of file paths to process

        Returns:
            List of processed document results
        """
        results = []
        self.logger.info(f"Processing batch of {len(file_paths)} documents")

        for i, file_path in enumerate(file_paths):
            try:
                result = self.process_document(file_path)
                results.append(result)
                self.logger.info(f"Processed {i+1}/{len(file_paths)}: {file_path}")
            except Exception as e:
                self.logger.error(f"❌ Failed to process {file_path}: {str(e)}")
                # Continue with other files
                continue

        return results

    def _extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        Extract common metadata from file.

        Args:
            file_path: Path to the file

        Returns:
            Dictionary containing file metadata
        """
        file_stat = os.stat(file_path)
        file_path_obj = Path(file_path)

        return {
            "filename": file_path_obj.name,
            "file_extension": file_path_obj.suffix.lower(),
            "file_size_bytes": file_stat.st_size,
            "file_size_mb": round(file_stat.st_size / (1024 * 1024), 2),
            "created_time": datetime.fromtimestamp(file_stat.st_ctime).isoformat(),
            "modified_time": datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
            "processed_time": datetime.now().isoformat(),
        }

    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        📄 Extract text from a PDF document using PyMuPDF with fallback to pdfplumber.

        Args:
            file_path: Path to the PDF file

        Returns:
            Dictionary with extracted text and metadata
        """
        self.logger.info(f"Processing PDF: {file_path}")

        text_content = []
        metadata = self._extract_metadata(file_path)

        try:
            # Primary method: PyMuPDF (faster)
            doc = fitz.open(file_path)
            metadata.update(
                {
                    "page_count": doc.page_count,
                    "title": doc.metadata.get("title", ""),
                    "author": doc.metadata.get("author", ""),
                    "subject": doc.metadata.get("subject", ""),
                    "creator": doc.metadata.get("creator", ""),
                }
            )

            for page_num in range(doc.page_count):
                page = doc[page_num]
                text = page.get_text()
                if text.strip():  # Only add non-empty pages
                    text_content.append({"page": page_num + 1, "content": text.strip()})

            doc.close()

        except Exception as e:
            self.logger.warning(f"PyMuPDF failed, trying pdfplumber: {str(e)}")

            # Fallback method: pdfplumber (more robust for complex PDFs)
            try:
                with pdfplumber.open(file_path) as pdf:
                    metadata["page_count"] = len(pdf.pages)

                    for page_num, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text and text.strip():
                            text_content.append(
                                {"page": page_num + 1, "content": text.strip()}
                            )

            except Exception as fallback_error:
                raise DocumentProcessingError(
                    f"Both PDF extraction methods failed: {str(fallback_error)}",
                    file_path,
                )

        # Final content processing
        full_text = "\n\n".join([item["content"] for item in text_content])
        metadata["total_characters"] = len(full_text)
        metadata["total_words"] = len(full_text.split())

        return {
            "content": full_text,
            "pages": text_content,
            "metadata": metadata,
            "source": file_path,
            "document_type": "pdf",
        }

    def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from a DOCX document using python-docx.

        Args:
            file_path: Path to the DOCX file

        Returns:
            Dictionary with extracted text and metadata
        """
        self.logger.info(f"Processing DOCX: {file_path}")

        try:
            doc = docx.Document(file_path)
            metadata = self._extract_metadata(file_path)

            # Extract document properties
            core_props = doc.core_properties
            metadata.update(
                {
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "subject": core_props.subject or "",
                    "created": (
                        core_props.created.isoformat() if core_props.created else ""
                    ),
                    "modified": (
                        core_props.modified.isoformat() if core_props.modified else ""
                    ),
                    "paragraph_count": len(doc.paragraphs),
                }
            )

            # Extract text content
            paragraphs = []
            full_text_parts = []

            for i, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if text:  # Only include non-empty paragraphs
                    paragraphs.append({"paragraph": i + 1, "content": text})
                    full_text_parts.append(text)

            #   Extract tables if present
            tables_content = []
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    if any(row_data):  # Only include non-empty rows
                        table_data.append(row_data)

                if table_data:
                    tables_content.append({"table": table_idx + 1, "data": table_data})
                    # Add table content to full text
                    table_text = "\n".join([" | ".join(row) for row in table_data])
                    full_text_parts.append(f"\n[Table {table_idx + 1}]\n{table_text}")

            full_text = "\n\n".join(full_text_parts)
            metadata.update(
                {
                    "total_characters": len(full_text),
                    "total_words": len(full_text.split()),
                    "table_count": len(tables_content),
                }
            )

            return {
                "content": full_text,
                "paragraphs": paragraphs,
                "tables": tables_content,
                "metadata": metadata,
                "source": file_path,
                "document_type": "docx",
            }

        except Exception as e:
            raise DocumentProcessingError(f"Error processing DOCX: {str(e)}", file_path)

    def _process_spreadsheet(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from a CSV or Excel file using pandas.

        Args:
            file_path: Path to the spreadsheet file

        Returns:
            Dictionary with extracted text and metadata
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        self.logger.info(f"Processing spreadsheet: {file_path}")

        try:
            metadata = self._extract_metadata(file_path)
            sheets_data = []

            if file_extension == ".csv":
                # 📄 Process CSV file
                df = pd.read_csv(file_path, encoding="utf-8")
                sheet_content = self._process_dataframe(df, "Sheet1")
                sheets_data.append(sheet_content)
                metadata["sheet_count"] = 1

            else:
                # Process Excel file
                excel_file = pd.ExcelFile(file_path)
                metadata["sheet_count"] = len(excel_file.sheet_names)

                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    sheet_content = self._process_dataframe(df, sheet_name)
                    sheets_data.append(sheet_content)

            # 🔗 Combine all sheets content
            full_text_parts = []
            for sheet in sheets_data:
                full_text_parts.append(f"[{sheet['sheet_name']}]\n{sheet['content']}")

            full_text = "\n\n".join(full_text_parts)
            metadata.update(
                {
                    "total_characters": len(full_text),
                    "total_words": len(full_text.split()),
                    "total_rows": sum(sheet["row_count"] for sheet in sheets_data),
                    "total_columns": (
                        max(sheet["column_count"] for sheet in sheets_data)
                        if sheets_data
                        else 0
                    ),
                }
            )

            return {
                "content": full_text,
                "sheets": sheets_data,
                "metadata": metadata,
                "source": file_path,
                "document_type": "spreadsheet",
            }

        except Exception as e:
            raise DocumentProcessingError(
                f"Error processing spreadsheet: {str(e)}", file_path
            )

    def _process_dataframe(self, df: pd.DataFrame, sheet_name: str) -> Dict[str, Any]:
        """
        Process a pandas DataFrame into text content.

        Args:
            df: Pandas DataFrame
            sheet_name: Name of the sheet

        Returns:
            Dictionary with processed sheet data
        """
        # Clean the dataframe
        df = df.dropna(how="all")  # Remove completely empty rows
        df = df.fillna("")  # Fill NaN with empty strings

        #   Create text representation
        content_parts = []

        # Add headers
        headers = df.columns.tolist()
        content_parts.append(" | ".join(str(h) for h in headers))
        content_parts.append("-" * 50)  # Separator

        # Add data rows
        for _, row in df.iterrows():
            row_text = " | ".join(str(cell) for cell in row.values)
            content_parts.append(row_text)

        content = "\n".join(content_parts)

        return {
            "sheet_name": sheet_name,
            "content": content,
            "headers": headers,
            "row_count": len(df),
            "column_count": len(df.columns),
            "data": df.to_dict("records"),  # For structured access
        }

    def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """
        🎯 Extract text from a PowerPoint presentation using python-pptx.

        Args:
            file_path: Path to the PPTX file

        Returns:
            Dictionary with extracted text and metadata
        """
        self.logger.info(f" Processing PPTX: {file_path}")

        try:
            presentation = pptx.Presentation(file_path)
            metadata = self._extract_metadata(file_path)

            # Extract presentation metadata
            core_props = presentation.core_properties
            metadata.update(
                {
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "subject": core_props.subject or "",
                    "created": (
                        core_props.created.isoformat() if core_props.created else ""
                    ),
                    "modified": (
                        core_props.modified.isoformat() if core_props.modified else ""
                    ),
                    "slide_count": len(presentation.slides),
                }
            )

            # 🎯 Extract content from slides
            slides_content = []
            full_text_parts = []

            for slide_idx, slide in enumerate(presentation.slides):
                slide_text_parts = []

                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text.strip())

                if slide_text_parts:
                    slide_content = "\n".join(slide_text_parts)
                    slides_content.append(
                        {"slide": slide_idx + 1, "content": slide_content}
                    )
                    full_text_parts.append(f"[Slide {slide_idx + 1}]\n{slide_content}")

            full_text = "\n\n".join(full_text_parts)
            metadata.update(
                {
                    "total_characters": len(full_text),
                    "total_words": len(full_text.split()),
                    "slides_with_content": len(slides_content),
                }
            )

            return {
                "content": full_text,
                "slides": slides_content,
                "metadata": metadata,
                "source": file_path,
                "document_type": "pptx",
            }

        except Exception as e:
            raise DocumentProcessingError(f"Error processing PPTX: {str(e)}", file_path)

    def _process_text_file(self, file_path: str) -> Dict[str, Any]:
        """
        📝 Extract text from plain text files (.txt, .md).

        Args:
            file_path: Path to the text file

        Returns:
            Dictionary with extracted text and metadata
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        self.logger.info(f" Processing text file: {file_path}")

        try:
            metadata = self._extract_metadata(file_path)

            # Try different encodings for robust text reading
            encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
            content = None

            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as file:
                        content = file.read()
                    self.logger.info(
                        f" Successfully read file with {encoding} encoding"
                    )
                    break
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    self.logger.warning(f"Failed to read with {encoding}: {str(e)}")
                    continue

            if content is None:
                raise DocumentProcessingError(
                    f"Could not read file with any supported encoding", file_path
                )

            # Clean and process content
            content = content.strip()
            if not content:
                raise DocumentProcessingError(
                    f"File is empty or contains no readable text", file_path
                )

            # Split content into logical sections for better processing
            sections = []
            if file_extension == ".md":
                # 📋 For Markdown files, split by headers
                sections = self._split_markdown_content(content)
            else:
                # 📄 For plain text, split by paragraphs
                sections = self._split_text_content(content)

            # Update metadata with text-specific information
            lines = content.split("\n")
            metadata.update(
                {
                    "file_type": (
                        "markdown" if file_extension == ".md" else "plain_text"
                    ),
                    "line_count": len(lines),
                    "paragraph_count": len(
                        [p for p in content.split("\n\n") if p.strip()]
                    ),
                    "total_characters": len(content),
                    "total_words": len(content.split()),
                    "encoding_used": encoding if "encoding" in locals() else "utf-8",
                    "sections_count": len(sections),
                }
            )

            return {
                "content": content,
                "sections": sections,
                "metadata": metadata,
                "source": file_path,
                "document_type": "markdown" if file_extension == ".md" else "text",
            }

        except Exception as e:
            raise DocumentProcessingError(
                f"Error processing text file: {str(e)}", file_path
            )

    def _split_markdown_content(self, content: str) -> List[Dict[str, Any]]:
        """
        Split Markdown content by headers for better organization.

        Args:
            content: Markdown content

        Returns:
            List of sections with headers and content
        """
        sections = []
        lines = content.split("\n")
        current_section = {"header": "", "content": [], "level": 0}

        for line in lines:
            # Check for markdown headers
            if line.strip().startswith("#"):
                # Save previous section if it has content
                if current_section["content"] or current_section["header"]:
                    section_content = "\n".join(current_section["content"]).strip()
                    if section_content or current_section["header"]:
                        sections.append(
                            {
                                "header": current_section["header"],
                                "content": section_content,
                                "level": current_section["level"],
                                "section_index": len(sections),
                            }
                        )

                # Start new section
                header_level = len(line) - len(line.lstrip("#"))
                header_text = line.lstrip("#").strip()
                current_section = {
                    "header": header_text,
                    "content": [],
                    "level": header_level,
                }
            else:
                current_section["content"].append(line)

        # Add the last section
        if current_section["content"] or current_section["header"]:
            section_content = "\n".join(current_section["content"]).strip()
            if section_content or current_section["header"]:
                sections.append(
                    {
                        "header": current_section["header"],
                        "content": section_content,
                        "level": current_section["level"],
                        "section_index": len(sections),
                    }
                )

        # If no headers found, treat entire content as one section
        if not sections:
            sections.append(
                {
                    "header": "Document Content",
                    "content": content.strip(),
                    "level": 1,
                    "section_index": 0,
                }
            )

        return sections

    def _split_text_content(self, content: str) -> List[Dict[str, Any]]:
        """
        Split plain text content by paragraphs.

        Args:
            content: Plain text content

        Returns:
            List of paragraph sections
        """
        sections = []
        paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]

        for i, paragraph in enumerate(paragraphs):
            sections.append(
                {
                    "header": f"Paragraph {i + 1}",
                    "content": paragraph,
                    "level": 1,
                    "section_index": i,
                }
            )

        # If no clear paragraphs, treat as single section
        if not sections:
            sections.append(
                {
                    "header": "Document Content",
                    "content": content.strip(),
                    "level": 1,
                    "section_index": 0,
                }
            )

        return sections
